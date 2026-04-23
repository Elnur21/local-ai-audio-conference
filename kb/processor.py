from pathlib import Path


def extract_text(filepath: str) -> str:
    ext = Path(filepath).suffix.lower()
    handlers = {
        ".pdf":  _pdf,
        ".docx": _docx,
        ".xlsx": _excel,
        ".xls":  _excel,
        ".pptx": _pptx,
        ".ppt":  _pptx,
        ".txt":  _txt,
    }
    if ext not in handlers:
        raise ValueError(f"Dəstəklənməyən format: {ext}")
    return handlers[ext](filepath)


def _pdf(path):
    import pdfplumber
    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                parts.append(t)
    return "\n".join(parts)


def _excel(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    rows = []
    for sheet in wb.worksheets:
        rows.append(f"[{sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            vals = [str(c) for c in row if c is not None]
            if vals:
                rows.append(" | ".join(vals))
    return "\n".join(rows)


def _docx(path):
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[Slayd {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    return "\n".join(lines)


def _txt(path):
    return Path(path).read_text(encoding="utf-8", errors="ignore")
